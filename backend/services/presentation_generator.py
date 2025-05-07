from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from typing import Dict
from datetime import datetime

class PresentationGenerator:
    def __init__(self):
        self.output_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'data', 'presentations')
        os.makedirs(self.output_dir, exist_ok=True)

    def create_presentation(self, content: Dict) -> str:
        """
        Create a PowerPoint presentation from the generated content.
        
        Args:
            content (Dict): Structured content from ContentGenerator
            
        Returns:
            str: Path to the generated presentation file
        """
        prs = Presentation()
        
        # Set slide dimensions to 16:9
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)

        # Add title slide
        self._add_title_slide(prs, content['title'], content['overview'])
        
        # Add content slides
        for slide_data in content['slides']:
            self._add_content_slide(prs, slide_data)
            
        # Add conclusion slide
        self._add_conclusion_slide(prs, content['conclusion'])

        # Save the presentation
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"presentation_{timestamp}.pptx"
        filepath = os.path.join(self.output_dir, filename)
        prs.save(filepath)
        
        return filepath

    def _add_title_slide(self, prs: Presentation, title: str, overview: str):
        """Add the title slide to the presentation."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Set subtitle/overview
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = overview
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)

    def _add_content_slide(self, prs: Presentation, slide_data: Dict):
        """Add a content slide to the presentation."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = slide_data['title']
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Add content
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        
        # Add main content
        p = tf.add_paragraph()
        p.text = slide_data['content']
        p.font.size = Pt(24)
        
        # Add key points
        for point in slide_data['key_points']:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(20)
            p.level = 1

    def _add_conclusion_slide(self, prs: Presentation, conclusion: str):
        """Add the conclusion slide to the presentation."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = "Conclusion"
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Add conclusion text
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        p = tf.add_paragraph()
        p.text = conclusion
        p.font.size = Pt(24) 