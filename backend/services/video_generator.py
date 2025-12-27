import os
import ffmpeg
import shutil  # Added for clearing temp directory
from datetime import datetime
from typing import Tuple, List
from pathlib import Path
import logging
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont, ImageEnhance
import io
import requests
import subprocess

logger = logging.getLogger(__name__)

# Use environment variables for FFmpeg paths, with fallback to system installation
FFMPEG_BIN = os.environ.get("FFMPEG_BINARY", "ffmpeg")
FFPROBE_BIN = os.environ.get("FFPROBE_BINARY", "ffprobe")

class VideoGenerator:
    def __init__(self):
        self.base_dir = Path(__file__).resolve().parent.parent
        self.output_dir = self.base_dir / 'data' / 'videos'
        self.temp_dir = self.base_dir / 'data' / 'temp'
        # Use relative path for logo instead of hardcoded Windows path
        self.logo_path = self.base_dir.parent / 'frontend' / 'assets' / 'logo.png'
        
        # Create necessary directories
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.temp_dir.mkdir(parents=True, exist_ok=True)
        
        # Check for FFmpeg installation
        self._check_ffmpeg_installation()
        
        logger.info(f"VideoGenerator initialized with output_dir: {self.output_dir}")
        logger.info(f"Temp directory: {self.temp_dir}")

    def _check_ffmpeg_installation(self):
        """Check if FFmpeg is installed and accessible."""
        try:
            result = subprocess.run([FFMPEG_BIN, "-version"], 
                                  capture_output=True, text=True, timeout=10)
            ffmpeg_info = (
                result.stdout.split('\n')[0] 
                if result.stdout 
                else 'FFmpeg version info'
            )
            logger.info(f"FFmpeg found: {ffmpeg_info}")
        except (subprocess.TimeoutExpired, FileNotFoundError):
            error_msg = (
                f"FFmpeg binary not accessible: {FFMPEG_BIN}. "
                "Please ensure ffmpeg is installed and accessible in the environment."
            )
            logger.error(error_msg)
            raise RuntimeError(error_msg)

    def create_video(self, presentation_path: str, audio_path: str) -> str:
        """
        Create a video by combining presentation slides with audio.
        
        Args:
            presentation_path (str): Path to the PowerPoint presentation
            audio_path (str): Path to the audio narration file
            
        Returns:
            str: Path to the generated video file
        """
        try:
            # Clear temp directory to prevent image mixups between videos
            if self.temp_dir.exists():
                shutil.rmtree(self.temp_dir)
            self.temp_dir.mkdir(parents=True, exist_ok=True)

            # Convert paths to absolute paths if they're relative
            presentation_path = self._ensure_absolute_path(presentation_path)
            audio_path = self._ensure_absolute_path(audio_path)
            
            # Verify files exist
            if not os.path.exists(presentation_path):
                raise FileNotFoundError(f"Presentation file not found: {presentation_path}")
            if not os.path.exists(audio_path):
                raise FileNotFoundError(f"Audio file not found: {audio_path}")
            
            logger.info(f"Creating video from presentation: {presentation_path}")
            logger.info(f"Using audio file: {audio_path}")
            
            # Convert presentation to images
            image_paths = self._convert_pptx_to_images(presentation_path)
            logger.info(f"Converted presentation to {len(image_paths)} images")
            
            # Get audio duration
            audio_duration = self._get_audio_duration(audio_path)
            logger.info(f"Audio duration: {audio_duration} seconds")
            
            # Calculate slide durations
            slide_durations = self._calculate_slide_durations(len(image_paths), audio_duration)
            logger.info(f"Slide durations: {slide_durations}")
            
            # Create video
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            output_path = self.output_dir / f"video_{timestamp}.mp4"
            
            self._create_video_from_images(image_paths, slide_durations, audio_path, str(output_path))
            logger.info(f"Video created successfully at: {output_path}")
            
            return str(output_path)
            
        except Exception as e:
            logger.error(f"Error creating video: {str(e)}")
            raise

    def _ensure_absolute_path(self, path: str) -> str:
        """Convert relative path to absolute path if needed."""
        if not os.path.isabs(path):
            return str(self.base_dir / path)
        return path

    def _convert_pptx_to_images(self, pptx_path: str) -> List[str]:
        """Convert PowerPoint slides to images with Unsplash backgrounds and styled text."""
        import textwrap  # Added for text wrapping
        
        UNSPLASH_ACCESS_KEY = "ZTZxqGOEqUt7juPx84I3SuRS6eqBmCGRKhOMRv15vTo"
        unsplash_url = "https://api.unsplash.com/photos/random"

        def fetch_unsplash_image(query, slide_idx):
            # Make Unsplash cache prompt-specific to prevent image mixups
            # Remove invalid characters for Windows file systems (<>:"/\|?*)
            safe_prompt = query.strip().lower().replace(" ", "_").replace(":", "_").replace('"', '_').replace('<', '_').replace('>', '_').replace('|', '_').replace('?', '_').replace('*', '_')[:50]
            cache_dir = self.temp_dir / f"unsplash_{safe_prompt}"
            cache_dir.mkdir(exist_ok=True)
            
            cache_path = cache_dir / f"slide_{slide_idx+1}.jpg"
            if cache_path.exists():
                return str(cache_path)
            params = {
                "query": query,
                "orientation": "landscape",
                "client_id": UNSPLASH_ACCESS_KEY
            }
            try:
                resp = requests.get(unsplash_url, params=params, timeout=10)
                if resp.status_code == 200:
                    img_url = resp.json()["urls"]["regular"]
                    img_data = requests.get(img_url).content
                    with open(cache_path, "wb") as f:
                        f.write(img_data)
                    return str(cache_path)
            except Exception as e:
                logger.warning(f"Failed to fetch Unsplash image for '{query}': {e}")
            return None

        def wrap_text(text, font, max_width):
            """Wrap text to fit within the given width."""
            if not text:
                return []
            
            lines = []
            paragraphs = text.split('\n')
            
            for paragraph in paragraphs:
                # Try to wrap the text using textwrap
                words = paragraph.split(' ')
                current_line = ""
                
                for word in words:
                    test_line = current_line + word + " "
                    
                    # Get text width using the font
                    try:
                        # Use ImageDraw to measure text width
                        temp_img = Image.new('RGB', (1, 1))
                        temp_draw = ImageDraw.Draw(temp_img)
                        bbox = temp_draw.textbbox((0, 0), test_line, font=font)
                        text_width = bbox[2] - bbox[0]
                    except:
                        # Fallback if font measurement fails
                        text_width = len(test_line) * 10  # Rough estimate
                    
                    if text_width <= max_width:
                        current_line = test_line
                    else:
                        if current_line:
                            lines.append(current_line.strip())
                        current_line = word + " "
                
                if current_line:
                    lines.append(current_line.strip())
            
            return lines

        try:
            logger.info(f"Converting PowerPoint to images: {pptx_path}")
            prs = Presentation(pptx_path)
            image_paths = []
            
            # Force specific pixel sizes for consistent rendering across environments
            # Use larger font sizes that will render consistently on both local and cloud
            try:
                # Try to use arial.ttf first, then fallback to other fonts, and finally default
                font = ImageFont.truetype("arial.ttf", 72)  # Larger size for better visibility
                small_font = ImageFont.truetype("arial.ttf", 48)  # Larger size for better visibility
                slide_num_font = ImageFont.truetype("arial.ttf", 32)  # Larger size for better visibility
            except Exception:
                try:
                    # Try alternative font paths that might be available in cloud environments
                    font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 72)
                    small_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 48)
                    slide_num_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 32)
                except Exception:
                    # Final fallback to default font with size approximation
                    logger.warning("Could not load specific fonts, using default fonts")
                    # Use ImageFont.load_default() which should be available everywhere
                    font = ImageFont.load_default()
                    small_font = ImageFont.load_default()
                    slide_num_font = ImageFont.load_default()

            # Extract topic from first slide title
            topic = None
            for shape in prs.slides[0].shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    topic = shape.text.strip()
                    break

            for i, slide in enumerate(prs.slides):
                image_path = self.temp_dir / f"slide_{i+1}.png"
                
                # Extract title and content from slide shapes
                title = None
                content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if not title:
                            title = shape.text.strip()
                        else:
                            content.append(shape.text.strip())

                # Generate relevant search query for background image
                search_query = f"{topic} {title}" if topic and title else title or "education"
                
                # Fetch Unsplash image
                img_file = fetch_unsplash_image(search_query, i)
                if img_file:
                    bg_img = Image.open(img_file).convert("RGB").resize((1920, 1080))
                    # Darken the image for text readability
                    enhancer = ImageEnhance.Brightness(bg_img)
                    bg_img = enhancer.enhance(0.6)
                else:
                    bg_img = Image.new('RGB', (1920, 1080), color='#f0f4ff')
                
                img = bg_img.copy()
                draw = ImageDraw.Draw(img)
                
                # Draw a semi-transparent overlay for text
                overlay = Image.new('RGBA', img.size, (255,255,255,120))
                img = Image.alpha_composite(img.convert('RGBA'), overlay)
                draw = ImageDraw.Draw(img)
                
                # Draw title (large, bold, colored) with text wrapping
                y = 80
                if title:
                    # Wrap title text to fit within slide width
                    title_lines = wrap_text(title, font, 1700)  # Leave some margin
                    for line in title_lines:
                        draw.text((80, y), line, font=font, fill='#4F8EF7')
                        y += 90  # Adjust spacing based on font size
                
                # Draw content lines (normal, dark) with text wrapping
                for content_item in content:
                    content_lines = wrap_text(content_item, small_font, 1700)  # Leave some margin
                    for line in content_lines:
                        draw.text((100, y), line, font=small_font, fill='#222')
                        y += 55  # Adjust spacing based on font size
                        
                        # Check if we're approaching the bottom of the slide
                        if y > 900:  # Leave space for slide number
                            logger.warning(f"Content approaching bottom of slide {i+1}, truncating...")
                            break
                    if y > 900:
                        break
                
                # Draw slide number (bottom right)
                slide_num = f"Slide {i+1}"
                draw.text((1700, 1000), slide_num, font=slide_num_font, fill='#888')
                
                img = img.convert('RGB')
                img.save(str(image_path))
                image_paths.append(str(image_path))
                logger.debug(f"Created image for slide {i+1}: {image_path}")
            
            return image_paths
        except Exception as e:
            logger.error(f"Error converting PowerPoint to images: {str(e)}")
            raise

    def _get_audio_duration(self, audio_path: str) -> float:
        """Get the duration of the audio file in seconds."""
        try:
            logger.info(f"Getting audio duration for: {audio_path}")
            logger.info(f"Using ffprobe binary: {FFPROBE_BIN}")
            logger.info(f"Audio file exists: {os.path.exists(audio_path)} at {audio_path}")
            if not os.path.exists(audio_path):
                raise FileNotFoundError(f"Audio file not found: {audio_path}")
                
            probe = ffmpeg.probe(audio_path, cmd=FFPROBE_BIN)
            duration = float(probe['format']['duration'])
            logger.info(f"Audio duration: {duration} seconds")
            return duration
        except ffmpeg.Error as e:
            logger.error(f"FFmpeg error getting audio duration: {str(e)}")
            raise
        except Exception as e:
            logger.error(f"Error getting audio duration: {str(e)}")
            raise

    def _calculate_slide_durations(self, num_slides: int, total_duration: float) -> List[float]:
        """Calculate how long each slide should be shown."""
        duration_per_slide = total_duration / num_slides
        return [duration_per_slide] * num_slides

    def _create_video_from_images(self, image_paths: List[str], durations: List[float], 
                                audio_path: str, output_path: str):
        """Create video from images and audio using FFmpeg, with Adura AI watermark."""
        try:
            logger.info(f"Creating video from {len(image_paths)} images")
            # Verify all files exist
            for img_path in image_paths:
                if not os.path.exists(img_path):
                    logger.error(f"Image file not found: {img_path}")
                    raise FileNotFoundError(f"Image file not found: {img_path}")
            if not os.path.exists(audio_path):
                logger.error(f"Audio file not found: {audio_path}")
                raise FileNotFoundError(f"Audio file not found: {audio_path}")
            if not self.logo_path.exists():
                logger.error(f"Logo file not found: {self.logo_path}")
                raise FileNotFoundError(f"Logo file not found: {self.logo_path}")

            # Create a temporary file with the list of images and their durations
            concat_file = self.temp_dir / "concat.txt"
            with open(concat_file, 'w') as f:
                for img_path, duration in zip(image_paths, durations):
                    f.write(f"file '{img_path}'\n")
                    f.write(f"duration {duration}\n")
            logger.info(f"Concat file created at: {concat_file}")
            with open(concat_file, 'r') as f:
                logger.info(f"Concat file contents:\n{f.read()}")
            
            # Run FFmpeg to create video with logo watermark
            cmd = [
                FFMPEG_BIN,
                '-f', 'concat',
                '-safe', '0',
                '-i', str(concat_file),
                '-i', audio_path,
                '-i', str(self.logo_path),
                '-filter_complex', '[2:v]scale=200:-1[logo];[0:v][logo]overlay=50:main_h-overlay_h-50:format=auto,format=yuv420p[v]',
                '-map', '[v]',
                '-map', '1:a',
                '-c:v', 'libx264',
                '-c:a', 'aac',
                '-b:a', '192k',
                '-shortest',
                '-y',
                output_path
            ]
            
            logger.info(f"Running FFmpeg command: {' '.join(cmd)}")
            
            try:
                subprocess.run(cmd, check=True, capture_output=True, text=True)
                logger.info(f"Video created successfully at: {output_path}")
                return output_path
            except subprocess.CalledProcessError as e:
                logger.error(f"FFmpeg error: {e.stderr}")
                raise Exception(f"ffmpeg error: {e.stderr}")

        except ffmpeg.Error as e:
            logger.error(f"FFmpeg error: {str(e)}")
            raise
        except Exception as e:
            logger.error(f"Error creating video: {str(e)}")
            raise 